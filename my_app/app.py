import os
import json
import re
import base64
import hashlib
import html
from datetime import datetime
import certifi

import pandas as pd
import requests
import streamlit as st
import streamlit.components.v1 as components
from docx import Document
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation
from pymongo import MongoClient
from pymongo.server_api import ServerApi

# ---------------------------------
# 기본 설정
# ---------------------------------
st.set_page_config(
    page_title="Hazel",
    page_icon="🤖",
    layout="wide"
)

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Crimson+Pro:wght@400;500;600&family=Inter:wght@400;500;600&display=swap');

/* ── 전체 배경 ── */
.stApp {
    background-color: #f5f0e8;
    font-family: 'Inter', sans-serif;
}
.block-container {
    padding-top: 1.5rem;
    max-width: 1100px !important;
}

/* ── 사이드바 ── */
section[data-testid="stSidebar"] {
    width: 340px !important;
    background-color: #ede8df !important;
    border-right: 1px solid #d9d3c7;
}

/* 모든 사이드바 버튼: 흰 배경 + 테두리로 경계 명확하게 */
section[data-testid="stSidebar"] .stButton > button {
    background-color: #fdf9f4 !important;
    border: 1.5px solid #ddd6cc !important;
    color: #3d3529 !important;
    text-align: left;
    border-radius: 10px !important;
    font-size: 0.88rem;
    font-weight: 500 !important;
    padding: 0.5rem 0.75rem;
    width: 100%;
    transition: all 0.15s ease;
}
section[data-testid="stSidebar"] .stButton > button *,
section[data-testid="stSidebar"] .stButton > button p {
    color: #3d3529 !important;
    font-weight: 500 !important;
}
section[data-testid="stSidebar"] .stButton > button:hover {
    background-color: #f0e8de !important;
    border-color: #c17f3e !important;
    color: #2c2416 !important;
}

/* 로그아웃 버튼: 살짝 강조 */
section[data-testid="stSidebar"] .stButton > button[kind="secondary"]:first-of-type {
    border-color: #b09070 !important;
}
section[data-testid="stSidebar"] h1,
section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3 {
    color: #2c2416 !important;
    font-size: 0.78rem !important;
    font-weight: 700 !important;
    letter-spacing: 0.08em !important;
    text-transform: uppercase;
}
section[data-testid="stSidebar"] label,
section[data-testid="stSidebar"] p,
section[data-testid="stSidebar"] span {
    color: #3d3529 !important;
}

/* ── 채팅 타이틀 ── */
.chat-title {
    font-family: 'Crimson Pro', Georgia, serif;
    font-size: clamp(1.6rem, 2.5vw, 2.2rem);
    font-weight: 600;
    color: #2c2416;
    margin-bottom: 0.5rem;
    line-height: 2.0;
    letter-spacing: -0.01em;
}

/* ── 사용자 메시지 (화면 오른쪽) ── */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) {
    flex-direction: row-reverse !important;
    margin-left: auto !important;
    margin-right: 0 !important;
    max-width: 78% !important;
    background-color: #c17f3e !important;
    border-radius: 20px 4px 20px 20px !important;
    border: none !important;
    box-shadow: 0 2px 8px rgba(193,127,62,0.28) !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) p,
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) li,
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) span,
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) div {
    color: #ffffff !important;
}

/* ── AI 메시지 (화면 왼쪽) ── */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) {
    margin-right: auto !important;
    margin-left: 0 !important;
    max-width: 85% !important;
    background-color: #ffffff !important;
    border-radius: 4px 20px 20px 20px !important;
    border: 1px solid #e8e2d9 !important;
    box-shadow: 0 2px 10px rgba(44,36,22,0.07) !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) p,
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) li,
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) span {
    color: #3d3529 !important;
}

/* ── 하단 채팅 바 컨테이너 ── */
[data-testid="stBottom"] {
    background: linear-gradient(to top, #f5f0e8 75%, rgba(245,240,232,0)) !important;
    padding-top: 16px !important;
    border-top: none !important;
}
[data-testid="stBottom"] > div {
    background: transparent !important;
}

/* ── 채팅 입력창 ── */
.stChatInput textarea {
    background-color: #ffffff !important;
    border: 1.5px solid #c9c2b6 !important;
    border-radius: 16px !important;
    color: #2c2416 !important;
    font-size: 0.95rem !important;
    box-shadow: 0 2px 8px rgba(44,36,22,0.06) !important;
}
.stChatInput textarea:focus {
    border-color: #c17f3e !important;
    box-shadow: 0 0 0 3px rgba(193,127,62,0.12) !important;
}

/* ── 검색 결과 카드 ── */
.result-card {
    padding: 14px 16px;
    border: 1px solid #ddd6cc;
    border-radius: 14px;
    background: #ffffff;
    margin-bottom: 10px;
    box-shadow: 0 2px 10px rgba(44,36,22,0.05);
}
.result-title {
    font-weight: 600;
    color: #2c2416;
    margin-bottom: 6px;
    font-size: 0.97rem;
}
.result-meta {
    color: #6b5e4e;
    font-size: 0.93rem;
    line-height: 1.55;
}
.result-meta a {
    color: #c17f3e;
    text-decoration: none;
}
.result-meta a:hover { text-decoration: underline; }

/* ── 검색 요약 배지 ── */
.search-summary {
    padding: 12px 14px;
    border-radius: 14px;
    border: 1px solid #ddd6cc;
    background: #faf7f3;
    margin-bottom: 12px;
}
.search-badge {
    display: inline-block;
    padding: 0.2rem 0.55rem;
    border-radius: 999px;
    background: #f0e8da;
    color: #7a5c35;
    font-size: 0.8rem;
    font-weight: 600;
    margin-right: 0.35rem;
    margin-bottom: 0.35rem;
}

/* ── 이미지 카드 ── */
.image-card {
    border: 1px solid #ddd6cc;
    border-radius: 18px;
    overflow: hidden;
    background: #ffffff;
    box-shadow: 0 4px 16px rgba(44,36,22,0.07);
    margin-bottom: 14px;
}
.image-card-meta { padding: 12px 14px 14px; }
.image-card-title {
    font-weight: 600;
    color: #2c2416;
    margin-bottom: 4px;
    line-height: 1.4;
}
.image-card-sub {
    color: #8a7560;
    font-size: 0.88rem;
    margin-bottom: 8px;
}

/* ── 미리보기 ── */
.preview-wrap {
    border: 1px solid #ddd6cc;
    border-radius: 12px;
    padding: 10px;
    background: #faf7f3;
    margin-top: 8px;
}

/* ── 공통 텍스트 ── */
h1, h2, h3, h4 { color: #2c2416 !important; }
p, li { color: #3d3529; }

/* ── 사이드바 토글(햄버거) 항상 표시 ── */
[data-testid="collapsedControl"],
[data-testid="stSidebarCollapsedControl"] {
    opacity: 1 !important;
    pointer-events: auto !important;
    visibility: visible !important;
}

/* ── 모바일: 툴바 강제 표시, 토글 버튼 숨김 ── */
@media (max-width: 768px) {
    [data-testid="stToolbar"] {
        opacity: 1 !important;
        pointer-events: auto !important;
    }
    #__toolbar_toggle_btn__ {
        display: none !important;
    }
}

/* ── Streamlit 툴바 숨김/표시 토글 ── */
body.hide-streamlit-toolbar [data-testid="stToolbar"],
body.hide-streamlit-toolbar [data-testid="stStatusWidget"],
body.hide-streamlit-toolbar [data-testid="stDecoration"],
body.hide-streamlit-toolbar .stAppToolbar {
    display: none !important;
}
</style>
""", unsafe_allow_html=True)

# 툴바 토글 버튼 (브라우저 localStorage에 상태 저장)
components.html(
    """
    <script>
      (function () {
        const STORAGE_KEY = "streamlit_toolbar_hidden";
        const doc = window.parent.document;
        if (!doc) return;

        const ensureButton = () => {
          let btn = doc.getElementById("toolbar-toggle-btn");
          if (!btn) {
            btn = doc.createElement("button");
            btn.id = "toolbar-toggle-btn";
            btn.type = "button";
            btn.style.position = "fixed";
            btn.style.top = "10px";
            btn.style.right = "10px";
            btn.style.zIndex = "99999";
            btn.style.border = "1px solid #d0c7bb";
            btn.style.background = "#fffaf3";
            btn.style.color = "#4a3c2d";
            btn.style.borderRadius = "999px";
            btn.style.padding = "6px 10px";
            btn.style.fontSize = "12px";
            btn.style.cursor = "pointer";
            btn.style.boxShadow = "0 2px 8px rgba(0,0,0,0.12)";
            doc.body.appendChild(btn);
          }
          return btn;
        };

        const applyState = (isHidden) => {
          doc.body.classList.toggle("hide-streamlit-toolbar", isHidden);
          const btn = ensureButton();
          btn.textContent = isHidden ? "툴바 보이기" : "툴바 숨기기";
        };

        const initialHidden = localStorage.getItem(STORAGE_KEY) === "1";
        applyState(initialHidden);

        const btn = ensureButton();
        btn.onclick = () => {
          const nextHidden = !doc.body.classList.contains("hide-streamlit-toolbar");
          localStorage.setItem(STORAGE_KEY, nextHidden ? "1" : "0");
          applyState(nextHidden);
        };
      })();
    </script>
    """,
    height=0,
    width=0,
)

# ---------------------------------
# MongoDB
# ---------------------------------
@st.cache_resource
def get_db():
    mongo_uri = os.getenv("MONGODB_URI")
    mongo_db_name = os.getenv("MONGODB_DB")

    if not mongo_uri:
        try:
            mongo_uri = st.secrets["MONGODB_URI"]
        except Exception:
            mongo_uri = None

    if not mongo_db_name:
        try:
            mongo_db_name = st.secrets["MONGODB_DB"]
        except Exception:
            mongo_db_name = "my_ai_chatbot_prod"

    if not mongo_uri:
        st.error("MONGODB_URI가 없습니다. 환경변수 또는 Streamlit secrets에 설정하세요.")
        st.stop()

    client = MongoClient(
        mongo_uri,
        server_api=ServerApi("1"),
        tls=True,
        tlsCAFile=certifi.where(),
        serverSelectionTimeoutMS=5000,
        connectTimeoutMS=20000,
        socketTimeoutMS=20000,
    )
    return client[mongo_db_name]

def get_chats_col():
    return get_db()["chats"]

def get_rag_col():
    return get_db()["rag_chunks"]

def init_mongo():
    try:
        get_chats_col().create_index([("username", 1), ("chat_id", 1)], unique=True)
        get_chats_col().create_index([("username", 1), ("updated_at", -1)])
    except Exception as e:
        st.warning(f"MongoDB 인덱스 생성 경고: {e}")

def init_rag_index():
    """Atlas Vector Search 인덱스 생성 (이미 존재하면 무시)."""
    try:
        col = get_rag_col()
        existing = list(col.list_search_indexes())
        names = [idx.get("name") for idx in existing]
        if "rag_vector_index" not in names:
            col.create_search_index({
                "name": "rag_vector_index",
                "type": "vectorSearch",
                "definition": {
                    "fields": [
                        {
                            "type": "vector",
                            "path": "embedding",
                            "numDimensions": 1536,
                            "similarity": "cosine"
                        },
                        {
                            "type": "filter",
                            "path": "username"
                        }
                    ]
                }
            })
    except Exception:
        pass

init_mongo()
init_rag_index()

# ---------------------------------
# OpenAI / NAVER
# ---------------------------------
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    try:
        api_key = st.secrets["OPENAI_API_KEY"]
    except Exception:
        api_key = None

if not api_key:
    st.error("OPENAI_API_KEY가 없습니다. 환경변수 또는 Streamlit secrets에 설정하세요.")
    st.stop()

try:
    NAVER_CLIENT_ID = st.secrets["NAVER_CLIENT_ID"]
    NAVER_CLIENT_SECRET = st.secrets["NAVER_CLIENT_SECRET"]
except Exception:
    NAVER_CLIENT_ID = os.getenv("NAVER_CLIENT_ID", "")
    NAVER_CLIENT_SECRET = os.getenv("NAVER_CLIENT_SECRET", "")

client = OpenAI(api_key=api_key)

# ---------------------------------
# 토큰 비용 설정 (per 1M tokens, USD)
# ---------------------------------
MODEL_PRICING = {
    "gpt-4o-mini":  {"input": 0.150,  "output": 0.600},
    "gpt-4.1-mini": {"input": 0.400,  "output": 1.600},
    "gpt-4.1":      {"input": 2.000,  "output": 8.000},
    "gpt-5.4":      {"input": 3.000,  "output": 15.000},
}
KRW_PER_USD = 1380

def calc_usage_display(model: str, input_tokens: int, output_tokens: int) -> str:
    pricing = MODEL_PRICING.get(model, {"input": 0, "output": 0})
    cost_usd = (input_tokens * pricing["input"] + output_tokens * pricing["output"]) / 1_000_000
    cost_krw = cost_usd * KRW_PER_USD
    total = input_tokens + output_tokens
    return (
        f"🔢 토큰: 입력 {input_tokens:,} · 출력 {output_tokens:,} · 합계 {total:,} &nbsp;|&nbsp; "
        f"💰 비용: ${cost_usd:.4f} (≈ ₩{cost_krw:.1f})"
    )

# ---------------------------------
# 로그인 관련
# ---------------------------------
def load_users():
    try:
        users = st.secrets.get("USERS", [])
        if isinstance(users, list) and len(users) > 0:
            return users
    except Exception:
        pass

    try:
        if os.path.exists("users.json"):
            with open("users.json", "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return data
    except Exception as e:
        st.warning(f"users.json 읽기 오류: {e}")

    return []

def verify_login(username: str, password: str) -> bool:
    users = load_users()
    username = str(username).strip()

    for user in users:
        if (
            str(user.get("username", "")).strip() == username
            and str(user.get("password", "")) == password
        ):
            return True
    return False

# ---------------------------------
# 파일 읽기 함수
# ---------------------------------
def read_pdf(file):
    try:
        file.seek(0)
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        return f"[PDF 읽기 실패: {e}]"

def read_excel(file):
    try:
        file.seek(0)
        excel = pd.ExcelFile(file)
        text_parts = []
        previews = []

        for sheet_name in excel.sheet_names:
            df = pd.read_excel(excel, sheet_name=sheet_name)
            previews.append((sheet_name, df.head(20)))
            text_parts.append(f"[시트: {sheet_name}]")
            text_parts.append(df.head(50).to_string(index=False))

        return "\n\n".join(text_parts), previews
    except Exception as e:
        return f"[Excel 읽기 실패: {e}]", []

def read_csv(file):
    try:
        file.seek(0)
        df = pd.read_csv(file)
        return df.head(50).to_string(index=False), df.head(20)
    except Exception:
        try:
            file.seek(0)
            df = pd.read_csv(file, encoding="cp949")
            return df.head(50).to_string(index=False), df.head(20)
        except Exception as e:
            return f"[CSV 읽기 실패: {e}]", None

def read_ppt(file):
    try:
        file.seek(0)
        prs = Presentation(file)
        text = ""
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    value = shape.text.strip()
                    if value:
                        slide_texts.append(value)
            text += f"\n[슬라이드 {i}]\n" + "\n".join(slide_texts) + "\n"
        return text.strip()
    except Exception as e:
        return f"[PPT 읽기 실패: {e}]"

def read_docx(file):
    try:
        file.seek(0)
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        return f"[DOCX 읽기 실패: {e}]"

def read_txt(file):
    file.seek(0)
    raw = file.getvalue()
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(enc)
        except Exception:
            pass
    return raw.decode("utf-8", errors="ignore")

def image_to_base64(file):
    file.seek(0)
    return base64.b64encode(file.getvalue()).decode()

# ---------------------------------
# RAG — 문서 임베딩 & 벡터 검색
# ---------------------------------
RAG_CHUNK_SIZE = 800
RAG_CHUNK_OVERLAP = 100
RAG_EMBED_MODEL = "text-embedding-3-small"
RAG_TOP_K = 5

def chunk_text(text: str, chunk_size: int = RAG_CHUNK_SIZE, overlap: int = RAG_CHUNK_OVERLAP) -> list[str]:
    """텍스트를 overlap이 있는 청크로 분할."""
    text = text.strip()
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end].strip())
        start += chunk_size - overlap
    return [c for c in chunks if c]

def embed_texts(texts: list[str]) -> list[list[float]]:
    """OpenAI 임베딩 API로 벡터 생성 (배치)."""
    texts = [t.replace("\n", " ") for t in texts]
    response = client.embeddings.create(model=RAG_EMBED_MODEL, input=texts)
    return [item.embedding for item in response.data]

def index_document(username: str, doc_id: str, filename: str, text: str) -> int:
    """문서를 청크로 분할 → 임베딩 → MongoDB 저장. 저장된 청크 수 반환."""
    chunks = chunk_text(text)
    if not chunks:
        return 0

    col = get_rag_col()
    col.delete_many({"username": username, "doc_id": doc_id})

    batch_size = 50
    total = 0
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        embeddings = embed_texts(batch)
        docs = [
            {
                "username": username,
                "doc_id": doc_id,
                "filename": filename,
                "chunk_index": i + j,
                "text": batch[j],
                "embedding": embeddings[j],
                "created_at": datetime.utcnow(),
            }
            for j in range(len(batch))
        ]
        col.insert_many(docs)
        total += len(docs)

    return total

def search_rag_chunks(username: str, query: str, top_k: int = RAG_TOP_K) -> list[dict]:
    """쿼리와 가장 유사한 청크를 벡터 검색으로 반환."""
    query_vec = embed_texts([query])[0]
    pipeline = [
        {
            "$vectorSearch": {
                "index": "rag_vector_index",
                "path": "embedding",
                "queryVector": query_vec,
                "numCandidates": top_k * 10,
                "limit": top_k,
                "filter": {"username": username},
            }
        },
        {
            "$project": {
                "_id": 0,
                "filename": 1,
                "chunk_index": 1,
                "text": 1,
                "score": {"$meta": "vectorSearchScore"},
            }
        },
    ]
    try:
        return list(get_rag_col().aggregate(pipeline))
    except Exception:
        return []

def list_rag_docs(username: str) -> list[dict]:
    """사용자의 RAG 문서 목록 반환 (doc_id + filename + 청크 수)."""
    pipeline = [
        {"$match": {"username": username}},
        {"$group": {
            "_id": "$doc_id",
            "filename": {"$first": "$filename"},
            "chunks": {"$sum": 1},
            "created_at": {"$first": "$created_at"},
        }},
        {"$sort": {"created_at": -1}},
    ]
    try:
        return [
            {"doc_id": d["_id"], "filename": d["filename"], "chunks": d["chunks"]}
            for d in get_rag_col().aggregate(pipeline)
        ]
    except Exception:
        return []

def delete_rag_doc(username: str, doc_id: str):
    """RAG 문서 삭제."""
    get_rag_col().delete_many({"username": username, "doc_id": doc_id})

# ---------------------------------
# 엑셀 변환 / 구조화 데이터 추출
# ---------------------------------
def dataframe_to_excel_bytes(df: pd.DataFrame, sheet_name="result"):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False, sheet_name=sheet_name)
    return output.getvalue()

def extract_json_block(text: str):
    if not text:
        return None

    patterns = [
        r"```json\s*(.*?)\s*```",
        r"```[\w]*\s*(\[\s*{.*?}\s*\])\s*```",
        r"```[\w]*\s*(\{\s*.*?\s*\})\s*```",
    ]

    for pattern in patterns:
        m = re.search(pattern, text, re.DOTALL)
        if m:
            candidate = m.group(1).strip()
            try:
                return json.loads(candidate)
            except Exception:
                pass

    stripped = text.strip()
    if stripped.startswith("[") or stripped.startswith("{"):
        try:
            return json.loads(stripped)
        except Exception:
            pass

    return None

def json_to_dataframe(data):
    if isinstance(data, list):
        if len(data) == 0:
            return pd.DataFrame()
        if all(isinstance(x, dict) for x in data):
            return pd.DataFrame(data)
        return pd.DataFrame({"value": data})

    if isinstance(data, dict):
        for value in data.values():
            if isinstance(value, list) and value and all(isinstance(x, dict) for x in value):
                return pd.DataFrame(value)
        return pd.DataFrame([data])

    return None

def parse_markdown_table_to_df(text: str):
    if not text or "|" not in text:
        return None

    lines = [line.strip() for line in text.splitlines() if "|" in line]
    if len(lines) < 2:
        return None

    header_idx = None
    for i in range(len(lines) - 1):
        if re.fullmatch(r"\|?[\s:-|]+\|?", lines[i + 1]):
            header_idx = i
            break

    if header_idx is None:
        return None

    header_line = lines[header_idx]
    data_lines = lines[header_idx + 2:]

    headers = [h.strip() for h in header_line.strip("|").split("|")]
    rows = []

    for line in data_lines:
        cells = [c.strip() for c in line.strip("|").split("|")]
        if len(cells) == len(headers):
            rows.append(cells)

    if not rows:
        return None

    return pd.DataFrame(rows, columns=headers)

def try_build_result_dataframe(full_text: str):
    data = extract_json_block(full_text)
    if data is not None:
        df = json_to_dataframe(data)
        if df is not None and not df.empty:
            return df

    df = parse_markdown_table_to_df(full_text)
    if df is not None and not df.empty:
        return df

    return None

# ---------------------------------
# HTML/CSS/JS 코드 추출 + 미리보기
# ---------------------------------
def extract_code_blocks(text: str):
    result = {
        "html": "",
        "css": "",
        "js": ""
    }

    if not text:
        return result

    matches = re.findall(r"```(\w+)?\s*(.*?)```", text, re.DOTALL)
    for lang, code in matches:
        lang = (lang or "").strip().lower()
        code = code.strip()

        if lang in ["html", "htm"]:
            result["html"] += "\n" + code
        elif lang == "css":
            result["css"] += "\n" + code
        elif lang in ["js", "javascript"]:
            result["js"] += "\n" + code

    return result

def build_preview_html_from_response(text: str):
    blocks = extract_code_blocks(text)

    html_code = blocks["html"].strip()
    css_code = blocks["css"].strip()
    js_code = blocks["js"].strip()

    if not html_code and not css_code and not js_code:
        return None, blocks

    if not html_code:
        return None, blocks

    if "<html" in html_code.lower():
        final_html = html_code

        if css_code:
            if "</head>" in final_html.lower():
                final_html = re.sub(
                    r"</head>",
                    f"<style>\n{css_code}\n</style>\n</head>",
                    final_html,
                    flags=re.IGNORECASE
                )
            else:
                final_html = f"<style>\n{css_code}\n</style>\n" + final_html

        if js_code:
            if "</body>" in final_html.lower():
                final_html = re.sub(
                    r"</body>",
                    f"<script>\n{js_code}\n</script>\n</body>",
                    final_html,
                    flags=re.IGNORECASE
                )
            else:
                final_html += f"\n<script>\n{js_code}\n</script>\n"

        return final_html, blocks

    final_html = f"""
<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8" />
<meta name="viewport" content="width=device-width, initial-scale=1.0" />
<style>
body {{
    font-family: Arial, sans-serif;
    padding: 20px;
    margin: 0;
    background: #ffffff;
}}
{css_code}
</style>
</head>
<body>
{html_code}
<script>
{js_code}
</script>
</body>
</html>
"""
    return final_html, blocks

def should_show_preview(user_input: str, response_text: str) -> bool:
    combined = f"{user_input}\n{response_text}".lower()

    keywords = [
        "html", "css", "js", "javascript",
        "퍼블리싱", "마크업", "웹페이지", "랜딩페이지",
        "코드", "미리보기", "화면 만들어", "ui 만들어"
    ]

    has_keyword = any(k in combined for k in keywords)
    has_html_block = "```html" in response_text.lower()

    return has_keyword or has_html_block

def should_prioritize_code_preview(query: str, files) -> bool:
    q = str(query).lower()
    code_keywords = [
        "html", "css", "js", "javascript",
        "시안", "퍼블리싱", "마크업", "웹페이지", "랜딩페이지",
        "코드", "화면", "ui", "ux"
    ]
    has_code_intent = any(k in q for k in code_keywords)

    has_image_attachment = False
    for f in files or []:
        name = str(getattr(f, "name", "")).lower()
        mime = str(getattr(f, "type", "")).lower()
        if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")) or mime.startswith("image/"):
            has_image_attachment = True
            break

    return has_code_intent and has_image_attachment

def render_preview_panel(preview_html: str, preview_blocks: dict, key_prefix: str):
    st.subheader("🖥 HTML/CSS 미리보기")

    if preview_blocks.get("html"):
        st.markdown("**HTML**")
        st.code(preview_blocks["html"], language="html")

    if preview_blocks.get("css"):
        st.markdown("**CSS**")
        st.code(preview_blocks["css"], language="css")

    if preview_blocks.get("js"):
        st.markdown("**JavaScript**")
        st.code(preview_blocks["js"], language="javascript")

    st.download_button(
        label="📥 코드 다운로드 (.html)",
        data=preview_html,
        file_name="preview.html",
        mime="text/html",
        key=f"{key_prefix}_download_html"
    )

    components.html(preview_html, height=700, scrolling=True)

# ---------------------------------
# 검색 유틸
# ---------------------------------
def safe_link(url: str) -> str:
    if not url:
        return ""
    return html.escape(url, quote=True)

def clean_html_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"<[^>]+>", "", text)
    return html.unescape(text).strip()

def should_search_web(query: str) -> bool:
    keywords = [
        "최신", "최근", "오늘", "뉴스", "이슈", "발표", "동향",
        "맛집", "카페", "식당", "술집", "브런치", "디저트",
        "근처", "어디", "추천", "여행", "여행지", "가볼만한 곳",
        "가격", "얼마", "출시", "일정", "오픈", "영업시간",
        "주가", "환율", "날씨", "후기", "리뷰", "순위",
        "뭐야", "왜", "어떻게", "정보", "찾아줘", "검색", "알려줘",
        "비교", "베스트", "인기", "사진", "이미지"
    ]
    return any(k in query for k in keywords)

def build_search_plan(query: str) -> dict:
    q = query.lower()

    local_keywords = [
        "맛집", "식당", "카페", "술집", "브런치", "디저트",
        "여행지", "가볼만한 곳", "데이트", "근처", "장소", "어디",
        "관광지", "볼거리", "놀거리", "숙소", "호텔"
    ]
    recency_keywords = [
        "뉴스", "최근", "최신", "오늘", "속보", "발표", "이슈",
        "논란", "동향", "주가", "환율", "날씨", "시세", "전망",
        "왜 떨어져", "왜 올랐", "무슨 일", "업데이트", "출시"
    ]
    image_keywords = [
        "사진", "이미지", "인테리어", "분위기", "외관", "메뉴", "비주얼"
    ]

    use_local = any(k in q for k in local_keywords)
    use_openai_web = any(k in q for k in recency_keywords)
    wants_images = any(k in q for k in image_keywords) or use_local

    use_naver_news = use_openai_web or "리뷰" in q or "후기" in q
    use_naver_web = use_local or use_openai_web or "추천" in q or "비교" in q

    mode_labels = []
    if use_local:
        mode_labels.append("네이버 로컬")
    if use_naver_news:
        mode_labels.append("네이버 뉴스")
    if use_naver_web:
        mode_labels.append("네이버 웹")
    if use_openai_web:
        mode_labels.append("OpenAI 웹검색")
    if wants_images:
        mode_labels.append("네이버 이미지")

    return {
        "use_local": use_local,
        "use_naver_news": use_naver_news,
        "use_naver_web": use_naver_web,
        "use_openai_web": use_openai_web,
        "wants_images": wants_images,
        "mode_labels": mode_labels or ["일반 응답"]
    }

def make_image_search_query(query: str) -> str:
    q = query.strip()
    if any(keyword in q for keyword in ["맛집", "식당", "카페", "브런치", "디저트"]):
        return f"{q} 분위기 메뉴"
    if any(keyword in q for keyword in ["여행", "여행지", "관광지", "호텔", "숙소"]):
        return f"{q} 사진"
    return q

def should_generate_image(query: str) -> bool:
    q = query.lower().strip()
    image_keywords = [
        "이미지", "사진", "그림", "일러스트", "포스터", "캐릭터", "배너", "썸네일", "로고",
        "메뉴", "메뉴판", "홍보물", "전단", "브로슈어", "쿠폰"
    ]
    action_keywords = ["만들어줘", "만들어 줘", "생성", "제작", "만들기", "뽑아줘"]
    direct_draw_keywords = [
        "그려줘", "그려 줘", "그림 그려", "일러스트", "포스터", "렌더링", "스케치",
        "디자인해줘", "디자인 해줘", "시안"
    ]
    return any(k in q for k in direct_draw_keywords) or (
        any(k in q for k in image_keywords) and any(k in q for k in action_keywords)
    )

def extract_image_generation_prompt(query: str) -> str:
    prompt = query.strip()
    cleanup_tokens = [
        "이미지 만들어줘", "이미지 생성해줘", "이미지 생성", "사진 만들어줘", "그림 그려줘",
        "그림 만들어줘", "일러스트 만들어줘", "포스터 만들어줘", "이미지로 만들어줘",
        "사진으로 만들어줘", "그려 줘", "그려줘"
    ]
    for token in cleanup_tokens:
        prompt = prompt.replace(token, "").strip()
    return prompt or query.strip()

def detect_image_generation_mode(query: str) -> str:
    q = str(query).lower()
    is_cafe = any(k in q for k in ["카페", "커피", "음료", "디저트"])
    wants_menu = any(k in q for k in ["메뉴", "메뉴판", "가격표"])
    wants_poster = any(k in q for k in ["포스터", "홍보", "신메뉴", "프로모션", "광고"])

    if is_cafe and wants_menu:
        return "cafe_menu"
    if is_cafe and wants_poster:
        return "cafe_poster"
    return "general"

def build_image_generation_spec(user_query: str, base_prompt: str) -> dict:
    mode = detect_image_generation_mode(user_query)
    if mode == "cafe_menu":
        enhanced_prompt = (
            f"{base_prompt}\n\n"
            "카페 메뉴판 디자인 스타일로 생성. 한국어 텍스트가 자연스럽고 읽기 쉬워야 함. "
            "섹션 구분(예: COFFEE / NON-COFFEE / TEA / ADE / DESSERT), 가격은 숫자로 선명하게 표기. "
            "배경 대비를 높여 가독성 우선, 메뉴명 정렬 깔끔하게. 군더더기 없는 상업용 메뉴판 톤."
        )
        return {"mode": mode, "prompt": enhanced_prompt, "size": "1536x1024", "quality": "high"}

    if mode == "cafe_poster":
        enhanced_prompt = (
            f"{base_prompt}\n\n"
            "카페 신메뉴 홍보 포스터 스타일로 생성. 세로형 레이아웃, 중앙 제품 히어로샷, "
            "상단에 임팩트 있는 헤드라인, 하단에 짧은 카피와 브랜드 무드. "
            "컬러는 청량하고 트렌디한 톤, 상업 광고물처럼 완성도 높게."
        )
        return {"mode": mode, "prompt": enhanced_prompt, "size": "1024x1536", "quality": "high"}

    return {"mode": mode, "prompt": base_prompt, "size": "1024x1024", "quality": "medium"}

# ---------------------------------
# 네이버 검색
# ---------------------------------
def naver_search(query: str, search_type: str = "local", display: int = 5):
    if not NAVER_CLIENT_ID or not NAVER_CLIENT_SECRET:
        return [{"error": "NAVER_CLIENT_ID / NAVER_CLIENT_SECRET 이 설정되지 않았습니다."}]

    url_map = {
        "webkr": "https://openapi.naver.com/v1/search/webkr.json",
        "news": "https://openapi.naver.com/v1/search/news.json",
        "blog": "https://openapi.naver.com/v1/search/blog.json",
        "local": "https://openapi.naver.com/v1/search/local.json",
    }

    url = url_map.get(search_type, url_map["local"])
    headers = {
        "X-Naver-Client-Id": NAVER_CLIENT_ID,
        "X-Naver-Client-Secret": NAVER_CLIENT_SECRET,
    }
    params = {
        "query": query,
        "display": display,
        "start": 1,
    }

    if search_type == "local":
        params["sort"] = "random"

    try:
        res = requests.get(url, headers=headers, params=params, timeout=15)
        res.raise_for_status()
        data = res.json()
        items = data.get("items", [])

        results = []
        for item in items:
            results.append({
                "title": clean_html_text(item.get("title", "")),
                "description": clean_html_text(item.get("description", "")),
                "link": item.get("link", ""),
                "originallink": item.get("originallink", ""),
                "roadAddress": item.get("roadAddress", ""),
                "address": item.get("address", ""),
                "category": item.get("category", ""),
                "telephone": item.get("telephone", ""),
            })

        if not results:
            return [{"error": "검색 결과가 없습니다."}]
        return results

    except Exception as e:
        return [{"error": f"네이버 검색 오류: {e}"}]

def naver_image_search(query: str, display: int = 6):
    if not NAVER_CLIENT_ID or not NAVER_CLIENT_SECRET:
        return [{"error": "NAVER_CLIENT_ID / NAVER_CLIENT_SECRET 이 설정되지 않았습니다."}]

    url = "https://openapi.naver.com/v1/search/image.json"
    headers = {
        "X-Naver-Client-Id": NAVER_CLIENT_ID,
        "X-Naver-Client-Secret": NAVER_CLIENT_SECRET,
    }
    params = {
        "query": query,
        "display": display,
        "start": 1,
        "sort": "sim"
    }

    try:
        res = requests.get(url, headers=headers, params=params, timeout=15)
        res.raise_for_status()
        data = res.json()
        items = data.get("items", [])

        results = []
        for item in items:
            results.append({
                "title": clean_html_text(item.get("title", "")),
                "link": item.get("link", ""),
                "thumbnail": item.get("thumbnail", ""),
                "sizeheight": item.get("sizeheight", ""),
                "sizewidth": item.get("sizewidth", "")
            })

        if not results:
            return [{"error": "이미지 검색 결과가 없습니다."}]
        return results

    except Exception as e:
        return [{"error": f"네이버 이미지 검색 오류: {e}"}]

def get_valid_image_results(image_results):
    return [
        item for item in (image_results or [])
        if "error" not in item and item.get("thumbnail")
    ]

def format_image_search_results(image_results) -> str:
    valid_items = get_valid_image_results(image_results)
    if not valid_items:
        return "이미지 검색 결과 없음"

    lines = []
    for i, item in enumerate(valid_items, start=1):
        lines.append(
            f"{i}. 제목: {item.get('title', '')}\n"
            f"   썸네일: {item.get('thumbnail', '')}\n"
            f"   원본 링크: {item.get('link', '')}"
        )
    return "\n\n".join(lines)

def generate_openai_image(prompt: str, size: str = "1024x1024", quality: str = "medium"):
    response = client.images.generate(
        model="gpt-image-1",
        prompt=prompt,
        size=size,
        quality=quality,
    )

    data = _to_dict(response).get("data", [])
    images = []
    for idx, item in enumerate(data, start=1):
        b64 = item.get("b64_json")
        image_url = item.get("url")
        image_bytes = None
        display_url = image_url

        if b64:
            image_bytes = base64.b64decode(b64)
            display_url = f"data:image/png;base64,{b64}"

        images.append({
            "id": idx,
            "prompt": prompt,
            "image_url": display_url,
            "image_bytes": image_bytes,
            "mime_type": "image/png"
        })

    return images

def format_search_summary(search_plan: dict, naver_results: dict, openai_web_sources: list) -> str:
    badges = "".join(
        f'<span class="search-badge">{html.escape(label)}</span>'
        for label in search_plan.get("mode_labels", [])
    )
    local_count = len([x for x in naver_results.get("local", []) if "error" not in x])
    news_count = len([x for x in naver_results.get("news", []) if "error" not in x])
    web_count = len([x for x in naver_results.get("web", []) if "error" not in x])
    image_count = len(get_valid_image_results(naver_results.get("image", [])))
    source_count = len(openai_web_sources)

    lines = [
        f"<div class='search-summary'><div><strong>검색 보강 모드</strong></div><div style='margin-top:8px'>{badges}</div>",
        f"<div style='margin-top:10px; color:#475569; font-size:0.92rem;'>네이버 로컬 {local_count}건 · 뉴스 {news_count}건 · 웹 {web_count}건 · 이미지 {image_count}건 · OpenAI 출처 {source_count}건</div></div>"
    ]
    return "".join(lines)

def format_naver_search_results(results, search_type="local") -> str:
    if not results:
        return "검색 결과 없음"

    lines = []
    for i, item in enumerate(results, start=1):
        if "error" in item:
            lines.append(f"{i}. 오류: {item['error']}")
            continue

        if search_type == "local":
            lines.append(
                f"{i}. 제목: {item.get('title','')}\n"
                f"   카테고리: {item.get('category','')}\n"
                f"   주소: {item.get('roadAddress') or item.get('address','')}\n"
                f"   전화: {item.get('telephone','')}\n"
                f"   링크: {item.get('link','')}"
            )
        else:
            lines.append(
                f"{i}. 제목: {item.get('title','')}\n"
                f"   요약: {item.get('description','')}\n"
                f"   링크: {item.get('originallink') or item.get('link','')}"
            )

    return "\n\n".join(lines)

def render_naver_search_results(results, search_type="local"):
    if not results:
        st.info("검색 결과가 없습니다.")
        return

    for i, item in enumerate(results, start=1):
        if "error" in item:
            st.warning(item["error"])
            continue

        if search_type == "local":
            link = safe_link(item.get("link", ""))
            telephone = html.escape(item.get("telephone", "") or "정보 없음")
            address = html.escape(item.get("roadAddress") or item.get("address", "") or "주소 정보 없음")
            category = html.escape(item.get("category", "") or "카테고리 없음")
            body = f"""
            <div class="result-card">
                <div class="result-title">{i}. {html.escape(item.get('title',''))}</div>
                <div class="result-meta">
                    <strong>카테고리</strong>: {category}<br>
                    <strong>주소</strong>: {address}<br>
                    <strong>전화</strong>: {telephone}<br>
                    {"<a href='" + link + "' target='_blank'>상세 링크 열기</a>" if link else ""}
                </div>
            </div>
            """
        else:
            raw_link = item.get("originallink") or item.get("link", "")
            link = safe_link(raw_link)
            description = html.escape(item.get("description", "") or "설명 없음")
            body = f"""
            <div class="result-card">
                <div class="result-title">{i}. {html.escape(item.get('title',''))}</div>
                <div class="result-meta">
                    요약: {description}<br>
                    {"<a href='" + link + "' target='_blank'>원문 열기</a>" if link else ""}
                </div>
            </div>
            """
        st.markdown(body, unsafe_allow_html=True)

def render_image_results(image_results):
    if not image_results:
        st.info("이미지 검색 결과가 없습니다.")
        return

    valid_items = get_valid_image_results(image_results)
    error_items = [item for item in image_results if "error" in item]

    for item in error_items:
        st.warning(item["error"])

    if not valid_items:
        return

    column_count = max(1, min(3, len(valid_items)))
    cols = st.columns(column_count)
    for idx, item in enumerate(valid_items):
        with cols[idx % len(cols)]:
            with st.container(border=True):
                _, center_col, _ = st.columns([1, 2, 1])
                with center_col:
                    st.image(item["thumbnail"], use_container_width=True)
                width = item.get("sizewidth") or "-"
                height = item.get("sizeheight") or "-"
                st.markdown(
                    f"""
                    <div class="image-card-meta">
                        <div class="image-card-title">{html.escape(item.get('title') or '이미지 결과')}</div>
                        <div class="image-card-sub">썸네일 크기: {html.escape(str(width))} × {html.escape(str(height))}</div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
                link = item.get("link", "")
                if link:
                    st.link_button("원본 보기", link, use_container_width=True)

def render_generated_images(generated_images):
    if not generated_images:
        return

    st.subheader("🎨 생성된 이미지")
    column_count = max(1, min(3, len(generated_images)))
    cols = st.columns(column_count)
    for idx, item in enumerate(generated_images):
        with cols[idx % len(cols)]:
            with st.container(border=True):
                _, center_col, _ = st.columns([1, 2, 1])
                with center_col:
                    st.image(item["image_url"], use_container_width=True)
                st.caption(item.get("prompt", "생성 이미지"))
                if item.get("image_bytes"):
                    st.download_button(
                        label=f"이미지 다운로드 {item['id']}",
                        data=item["image_bytes"],
                        file_name=f"generated_image_{item['id']}.png",
                        mime=item.get("mime_type", "image/png"),
                        key=f"download_generated_{item['id']}_{hash(item.get('prompt', ''))}"
                    )

# ---------------------------------
# OpenAI 웹검색
# ---------------------------------
def _to_dict(obj):
    if obj is None:
        return None
    if isinstance(obj, (str, int, float, bool)):
        return obj
    if isinstance(obj, list):
        return [_to_dict(x) for x in obj]
    if isinstance(obj, dict):
        return {k: _to_dict(v) for k, v in obj.items()}
    if hasattr(obj, "model_dump"):
        return _to_dict(obj.model_dump())
    if hasattr(obj, "dict"):
        return _to_dict(obj.dict())
    if hasattr(obj, "__dict__"):
        return _to_dict(vars(obj))
    return str(obj)

def extract_openai_web_sources(response) -> list:
    data = _to_dict(response)
    output_items = data.get("output", []) if isinstance(data, dict) else []

    sources = []
    seen = set()

    for item in output_items:
        if not isinstance(item, dict):
            continue
        if item.get("type") == "web_search_call":
            action = item.get("action", {}) or {}
            action_sources = action.get("sources", []) or []

            for src in action_sources:
                if not isinstance(src, dict):
                    continue
                url = src.get("url", "")
                src_type = src.get("type", "")
                if url and url not in seen:
                    seen.add(url)
                    sources.append({
                        "type": src_type,
                        "url": url
                    })

    return sources

def render_openai_web_sources(sources: list):
    if not sources:
        st.info("웹검색 출처가 없습니다.")
        return

    for i, src in enumerate(sources, start=1):
        url = src.get("url", "")
        src_type = src.get("type", "")

        st.markdown(
            f"""
            <div class="result-card">
                <div class="result-title">{i}. 출처</div>
                <div class="result-meta">
                    유형: {html.escape(src_type)}<br>
                    <a href="{html.escape(url, quote=True)}" target="_blank">링크 열기</a>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )

def _extract_usage(response) -> tuple[int, int]:
    """(input_tokens, output_tokens) 반환. 없으면 (0, 0)."""
    try:
        usage = getattr(response, "usage", None)
        if usage is None:
            return 0, 0
        inp = getattr(usage, "input_tokens", 0) or 0
        out = getattr(usage, "output_tokens", 0) or 0
        return int(inp), int(out)
    except Exception:
        return 0, 0

def run_openai_web_search(model_name: str, instructions: str, history_for_model: list, user_content: list):
    common_kwargs = dict(
        model=model_name,
        instructions=instructions,
        input=[
            *history_for_model,
            {"role": "user", "content": user_content}
        ],
        tool_choice="auto",
        include=["web_search_call.action.sources"]
    )

    try:
        response = client.responses.create(
            tools=[
                {
                    "type": "web_search",
                    "search_context_size": "medium"
                }
            ],
            **common_kwargs
        )
        return response.output_text, extract_openai_web_sources(response), _extract_usage(response)

    except Exception:
        response = client.responses.create(
            tools=[
                {
                    "type": "web_search_preview",
                    "search_context_size": "medium",
                    "search_content_types": ["text", "image"],
                    "user_location": {
                        "type": "approximate",
                        "city": "Seoul",
                        "country": "KR",
                        "region": "Seoul",
                        "timezone": "Asia/Seoul"
                    }
                }
            ],
            **common_kwargs
        )
        return response.output_text, extract_openai_web_sources(response), _extract_usage(response)

def is_image_generation_request(query: str) -> bool:
    matcher = globals().get("should_generate_image")
    if callable(matcher):
        return matcher(query)

    q = str(query).lower().strip()
    image_keywords = [
        "이미지", "사진", "그림", "일러스트", "포스터", "캐릭터", "배너", "썸네일", "로고",
        "메뉴", "메뉴판", "홍보물", "전단", "브로슈어", "쿠폰"
    ]
    action_keywords = ["만들어줘", "만들어 줘", "생성", "제작", "만들기", "뽑아줘"]
    direct_draw_keywords = [
        "그려줘", "그려 줘", "그림 그려", "일러스트", "포스터", "렌더링", "스케치",
        "디자인해줘", "디자인 해줘", "시안"
    ]
    return any(k in q for k in direct_draw_keywords) or (
        any(k in q for k in image_keywords) and any(k in q for k in action_keywords)
    )

def get_image_generation_prompt(query: str) -> str:
    extractor = globals().get("extract_image_generation_prompt")
    if callable(extractor):
        return extractor(query)
    return str(query).strip()

def get_default_runtime_state():
    return {
        "naver_results_map": {"local": [], "news": [], "web": [], "image": []},
        "openai_web_sources": [],
        "generated_images": [],
        "search_plan": {"mode_labels": ["일반 응답"]},
        "do_search": False,
    }

# ---------------------------------
# 대화 저장 함수 (MongoDB)
# ---------------------------------
def get_default_chat_data():
    return {
        "title": "새 대화",
        "agent_role": "",
        "messages": [
            {"role": "assistant", "content": "안녕하세요! 무엇을 도와드릴까요?"}
        ]
    }

def create_new_chat():
    chat_id = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    username = st.session_state.get("username", "guest")
    data = get_default_chat_data()

    get_chats_col().insert_one({
        "username": username,
        "chat_id": chat_id,
        "title": data["title"],
        "agent_role": data["agent_role"],
        "messages": data["messages"],
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow()
    })
    return chat_id

def load_chat(chat_id: str):
    username = st.session_state.get("username", "guest")
    doc = get_chats_col().find_one(
        {"username": username, "chat_id": chat_id},
        {"_id": 0, "title": 1, "agent_role": 1, "messages": 1}
    )

    if doc:
        return {
            "title": doc.get("title", "새 대화"),
            "agent_role": doc.get("agent_role", ""),
            "messages": doc.get("messages", get_default_chat_data()["messages"])
        }

    return get_default_chat_data()

def append_message(chat_id: str, role: str, content: str):
    username = st.session_state.get("username", "guest")

    get_chats_col().update_one(
        {"username": username, "chat_id": chat_id},
        {
            "$push": {
                "messages": {
                    "role": role,
                    "content": content
                }
            },
            "$set": {
                "updated_at": datetime.utcnow()
            },
            "$setOnInsert": {
                "created_at": datetime.utcnow(),
                "title": "새 대화"
            }
        },
        upsert=True
    )

def update_chat_title(chat_id: str, title: str):
    username = st.session_state.get("username", "guest")

    get_chats_col().update_one(
        {"username": username, "chat_id": chat_id},
        {
            "$set": {
                "title": title,
                "updated_at": datetime.utcnow()
            }
        }
    )

def update_chat_agent_role(chat_id: str, agent_role: str):
    username = st.session_state.get("username", "guest")
    get_chats_col().update_one(
        {"username": username, "chat_id": chat_id},
        {
            "$set": {
                "agent_role": agent_role.strip(),
                "updated_at": datetime.utcnow()
            }
        }
    )

def list_chats():
    username = st.session_state.get("username", "guest")

    docs = list(
        get_chats_col()
        .find(
            {"username": username},
            {"_id": 0, "chat_id": 1, "title": 1, "updated_at": 1}
        )
        .sort("updated_at", -1)
    )

    result = []
    for doc in docs:
        result.append({
            "id": doc["chat_id"],
            "title": doc.get("title", "제목 없음")
        })
    return result

def delete_chat(chat_id: str):
    username = st.session_state.get("username", "guest")
    get_chats_col().delete_one({"username": username, "chat_id": chat_id})

def search_chats(keyword: str):
    """키워드로 대화 제목 및 메시지 내용을 검색한다."""
    username = st.session_state.get("username", "guest")
    if not keyword.strip():
        return []

    regex = {"$regex": keyword.strip(), "$options": "i"}
    docs = list(
        get_chats_col()
        .find(
            {
                "username": username,
                "$or": [
                    {"title": regex},
                    {"messages.content": regex},
                ]
            },
            {"_id": 0, "chat_id": 1, "title": 1, "updated_at": 1, "messages": 1}
        )
        .sort("updated_at", -1)
        .limit(30)
    )

    results = []
    for doc in docs:
        snippets = []
        kw_lower = keyword.strip().lower()
        for msg in doc.get("messages", []):
            content = msg.get("content", "")
            if kw_lower in content.lower():
                idx = content.lower().find(kw_lower)
                start = max(0, idx - 25)
                end = min(len(content), idx + len(keyword) + 60)
                snippet = ("…" if start > 0 else "") + content[start:end].replace("\n", " ") + ("…" if end < len(content) else "")
                snippets.append(snippet)
                if len(snippets) >= 2:
                    break
        results.append({
            "id": doc["chat_id"],
            "title": doc.get("title", "제목 없음"),
            "snippets": snippets,
        })
    return results

def make_title_from_messages(messages):
    for msg in messages:
        if msg["role"] == "user":
            text = msg["content"].strip().replace("\n", " ")
            return text[:20] if len(text) > 20 else text
    return "새 대화"

def get_chat_avatar(role: str):
    if role == "user":
        return "🧸"
    return "🐰"

# ---------------------------------
# 프롬프트
# ---------------------------------
def build_system_prompt(answer_length: str, agent_role: str = "") -> str:
    if answer_length == "짧게":
        length_rule = "답변은 핵심만 5~7문장으로 간단히 설명한다."
    elif answer_length == "보통":
        length_rule = "답변은 8~12문장 정도로 설명한다."
    else:
        length_rule = "답변은 충분히 자세하게 설명하고, 필요하면 예시와 항목 정리를 포함한다."

    role_instruction = ""
    if agent_role and agent_role.strip():
        role_instruction = f"\n추가 역할 지시:\n- {agent_role.strip()}\n- 위 역할의 관점/톤을 유지하되, 사실을 지어내지는 않는다.\n"

    return f"""
너는 친절하고 유능한 한국어 AI 챗봇이다.
물어보는 언어에 맞게 대답한다.
모르는 내용은 추측하지 말고 불확실하다고 말한다.
사용자가 파일을 첨부한 경우 첨부 내용을 우선 참고한다.
사용자가 이미지(여권, 비자, 신분증, 계약서, 문서 캡처 등)를 첨부하면 이미지 자체를 직접 판독해서 답변한다.

이미지에서 특히 아래 정보가 있으면 정리한다.
- 이름
- 여권번호
- 국적
- 생년월일
- 발급일
- 만료일
- 비자 종류
- 체류기간

생년월일, 발급일, 만료일은 가능하면 YYYY-MM-DD 형태로 정리한다.
확실하지 않은 값은 추정이라고 표시하거나 비워둘 수 있다.
이미지 속 텍스트가 흐리거나 일부 가려져 있으면 보이는 범위 내에서만 답변한다.
사용자가 표, 엑셀, 리스트, 정리본을 요청하면 가능하면 JSON 배열 또는 표 형태로 구조화해서 제공한다.

사용자가 HTML/CSS/JS 코드 또는 웹 화면 마크업을 요청하면:
- 가능하면 반드시 ```html``` / ```css``` / ```javascript``` 코드블록으로 나누어 제공한다.
- HTML은 바로 브라우저에서 렌더 가능한 형태로 작성한다.
- CSS가 있으면 별도 ```css``` 블록으로 준다.
- 필요한 경우 간단한 JS도 ```javascript``` 블록으로 준다.

웹검색 결과가 함께 제공된 경우, 그 결과를 참고해서 답변하되 검색 결과에 없는 내용을 지어내지 않는다.
관련 이미지 검색 결과가 함께 제공된 경우, 시각적으로 참고할 수 있다고만 생각하고 사실관계는 텍스트 검색 결과를 우선한다.
{role_instruction}

{length_rule}
"""

# ---------------------------------
# 세션 초기화
# ---------------------------------
if "chat_paste_hint" not in st.session_state:
    st.session_state.chat_paste_hint = ""
    
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if "username" not in st.session_state:
    st.session_state.username = None

if "uploaded_files_cache" not in st.session_state:
    st.session_state.uploaded_files_cache = []

if "answer_length" not in st.session_state:
    st.session_state.answer_length = "보통"

if "model_name" not in st.session_state:
    st.session_state.model_name = "gpt-4.1-mini"

if "last_result_df" not in st.session_state:
    st.session_state.last_result_df = None

if "last_preview_html" not in st.session_state:
    st.session_state.last_preview_html = None

if "last_preview_blocks" not in st.session_state:
    st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}

if "use_web_search" not in st.session_state:
    st.session_state.use_web_search = True

if "auto_search_only" not in st.session_state:
    st.session_state.auto_search_only = True

if "show_search_images" not in st.session_state:
    st.session_state.show_search_images = True

if "chat_search_query" not in st.session_state:
    st.session_state.chat_search_query = ""

if "rag_enabled" not in st.session_state:
    st.session_state.rag_enabled = False

if "rag_indexing" not in st.session_state:
    st.session_state.rag_indexing = False

if "show_web_sources" not in st.session_state:
    st.session_state.show_web_sources = True

if "last_generated_images" not in st.session_state:
    st.session_state.last_generated_images = []

if "last_generated_prompt" not in st.session_state:
    st.session_state.last_generated_prompt = ""

if "is_generating" not in st.session_state:
    st.session_state.is_generating = False

if "stop_generation" not in st.session_state:
    st.session_state.stop_generation = False

if "last_paste_signature" not in st.session_state:
    st.session_state.last_paste_signature = ""

# ---------------------------------
# 로그인 화면
# ---------------------------------
if not st.session_state.logged_in:

    st.markdown("""
    <style>
    /* 로그인 전용 — 블록 컨테이너 풀스크린으로 */
    .block-container {
        max-width: 100% !important;
        padding: 0 !important;
    }

    /* ── 배경 블롭 ── */
    .login-blob {
        position: fixed;
        border-radius: 50%;
        pointer-events: none;
        z-index: 0;
        filter: blur(60px);
    }

    /* ── 카드 래퍼 ── */
    .login-card {
        background: rgba(255,255,255,0.82);
        backdrop-filter: blur(12px);
        -webkit-backdrop-filter: blur(12px);
        border: 1px solid rgba(212,201,184,0.7);
        border-radius: 28px;
        padding: 48px 44px 40px;
        box-shadow:
            0 4px 24px rgba(44,36,22,0.08),
            0 1px 4px rgba(44,36,22,0.06);
        margin-bottom: 1.5rem;
    }

    .login-logo {
        font-size: 3.2rem;
        line-height: 1;
        margin-bottom: 0.6rem;
        display: block;
        text-align: center;
    }
    .login-title {
        font-family: 'Crimson Pro', Georgia, serif;
        font-size: 2rem;
        font-weight: 600;
        color: #2c2416;
        text-align: center;
        margin-bottom: 0.3rem;
        letter-spacing: -0.01em;
    }
    .login-subtitle {
        font-size: 0.88rem;
        color: #8a7560;
        text-align: center;
        margin-bottom: 2.2rem;
    }
    .login-label {
        font-size: 0.82rem;
        font-weight: 600;
        color: #6b5e4e;
        margin-bottom: 0.3rem;
        display: block;
    }

    /* 입력 필드 */
    .stTextInput > div > div > input {
        border-radius: 14px !important;
        border: 1.5px solid #d4c9b8 !important;
        background: #faf7f3 !important;
        color: #2c2416 !important;
        padding: 0.65rem 1rem !important;
        font-size: 0.95rem !important;
    }
    .stTextInput > div > div > input:focus {
        border-color: #c17f3e !important;
        box-shadow: 0 0 0 3px rgba(193,127,62,0.14) !important;
        background: #ffffff !important;
    }

    /* 로그인 버튼 */
    .stButton > button {
        background: linear-gradient(135deg, #c17f3e 0%, #a66a2c 100%) !important;
        color: #ffffff !important;
        border: none !important;
        border-radius: 14px !important;
        font-size: 0.97rem !important;
        font-weight: 600 !important;
        padding: 0.7rem 1rem !important;
        letter-spacing: 0.02em !important;
        box-shadow: 0 4px 14px rgba(193,127,62,0.35) !important;
        transition: all 0.2s ease !important;
        margin-top: 0.6rem !important;
    }
    .stButton > button:hover {
        transform: translateY(-1px) !important;
        box-shadow: 0 6px 20px rgba(193,127,62,0.45) !important;
    }
    .stButton > button:active {
        transform: translateY(0px) !important;
    }

    /* divider 점 장식 */
    .login-dots {
        display: flex;
        justify-content: center;
        gap: 6px;
        margin-top: 1.8rem;
    }
    .login-dot {
        width: 6px; height: 6px;
        border-radius: 50%;
        background: #d4c9b8;
    }
    .login-dot.active { background: #c17f3e; }
    </style>

    <!-- 배경 블롭 3개 -->
    <div class="login-blob" style="
        width:520px; height:520px;
        top:-140px; right:-140px;
        background: radial-gradient(circle, rgba(193,127,62,0.22) 0%, transparent 70%);
    "></div>
    <div class="login-blob" style="
        width:420px; height:420px;
        bottom:-120px; left:-120px;
        background: radial-gradient(circle, rgba(44,36,22,0.13) 0%, transparent 70%);
    "></div>
    <div class="login-blob" style="
        width:300px; height:300px;
        top:55%; left:55%;
        background: radial-gradient(circle, rgba(193,127,62,0.10) 0%, transparent 70%);
    "></div>
    """, unsafe_allow_html=True)

    # 상단 여백
    st.markdown("<div style='height:60px'></div>", unsafe_allow_html=True)

    # 좌우 여백으로 중앙 카드
    left, center, right = st.columns([1, 1.1, 1])

    with center:
        st.markdown("""
        <div class="login-card">
            <span class="login-logo">🤖</span>
            <div class="login-title">Hazel</div>
            <div class="login-subtitle">계속하려면 로그인하세요</div>
        </div>
        """, unsafe_allow_html=True)

        login_username = st.text_input(
            "아이디",
            placeholder="아이디를 입력하세요",
        )
        login_password = st.text_input(
            "비밀번호",
            type="password",
            placeholder="비밀번호를 입력하세요",
        )

        if st.button("로그인", use_container_width=True):
            if verify_login(login_username, login_password):
                st.session_state.logged_in = True
                st.session_state.username = login_username
                st.rerun()
            else:
                st.error("아이디 또는 비밀번호가 올바르지 않습니다.")

        st.markdown("""
        <div class="login-dots">
            <div class="login-dot active"></div>
            <div class="login-dot"></div>
            <div class="login-dot"></div>
        </div>
        """, unsafe_allow_html=True)

    st.stop()

st.markdown('<div class="chat-title">🤖 Hazel</div>', unsafe_allow_html=True)

# ---------------------------------
# 로그인 후 현재 대화 초기화
# ---------------------------------
if "current_chat_id" not in st.session_state:
    chats = list_chats()
    if chats:
        st.session_state.current_chat_id = chats[0]["id"]
    else:
        st.session_state.current_chat_id = create_new_chat()

# ---------------------------------
# 사이드바
# ---------------------------------
with st.sidebar:
    st.write(f"로그인 사용자: **{st.session_state.username}**")

    if st.button("로그아웃", use_container_width=True):
        st.session_state.logged_in = False
        st.session_state.username = None
        st.session_state.uploaded_files_cache = []
        st.session_state.last_result_df = None
        st.session_state.last_preview_html = None
        st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}
        st.session_state.last_generated_images = []
        st.session_state.last_generated_prompt = ""
        st.session_state.last_paste_signature = ""
        if "current_chat_id" in st.session_state:
            del st.session_state["current_chat_id"]
        st.rerun()

    st.divider()
    st.header("대화")

    if st.button("＋ 새 대화", use_container_width=True):
        st.session_state.current_chat_id = create_new_chat()
        st.session_state.uploaded_files_cache = []
        st.session_state.last_result_df = None
        st.session_state.last_preview_html = None
        st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}
        st.session_state.last_generated_images = []
        st.session_state.last_generated_prompt = ""
        st.session_state.last_paste_signature = ""
        st.rerun()

    st.divider()

    search_input = st.text_input(
        "대화 검색",
        value=st.session_state.chat_search_query,
        placeholder="키워드로 검색...",
        key="sidebar_chat_search",
        label_visibility="collapsed",
    )
    if search_input != st.session_state.chat_search_query:
        st.session_state.chat_search_query = search_input
        st.rerun()

    if st.session_state.chat_search_query.strip():
        search_results = search_chats(st.session_state.chat_search_query)
        if search_results:
            st.caption(f"검색 결과 {len(search_results)}건")
            for res in search_results:
                col1, col2 = st.columns([4, 1])
                with col1:
                    if st.button(res["title"], key=f"srch_{res['id']}", use_container_width=True):
                        st.session_state.current_chat_id = res["id"]
                        st.session_state.chat_search_query = ""
                        st.session_state.uploaded_files_cache = []
                        st.session_state.last_result_df = None
                        st.session_state.last_preview_html = None
                        st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}
                        st.session_state.last_generated_images = []
                        st.session_state.last_generated_prompt = ""
                        st.session_state.last_paste_signature = ""
                        st.rerun()
                with col2:
                    if st.button("🗑", key=f"srch_del_{res['id']}", use_container_width=True):
                        deleting_current = (st.session_state.current_chat_id == res["id"])
                        delete_chat(res["id"])
                        remaining = list_chats()
                        if deleting_current:
                            if remaining:
                                st.session_state.current_chat_id = remaining[0]["id"]
                            else:
                                st.session_state.current_chat_id = create_new_chat()
                        st.rerun()
                for snippet in res["snippets"]:
                    st.caption(f"…{snippet}…")
        else:
            st.caption("검색 결과 없음")
    else:
        for chat in list_chats():
            col1, col2 = st.columns([4, 1])

            with col1:
                if st.button(chat["title"], key=f"open_{chat['id']}", use_container_width=True):
                    st.session_state.current_chat_id = chat["id"]
                    st.session_state.uploaded_files_cache = []
                    st.session_state.last_result_df = None
                    st.session_state.last_preview_html = None
                    st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}
                    st.session_state.last_generated_images = []
                    st.session_state.last_generated_prompt = ""
                    st.session_state.last_paste_signature = ""
                    st.rerun()

            with col2:
                if st.button("🗑", key=f"del_{chat['id']}", use_container_width=True):
                    deleting_current = (st.session_state.current_chat_id == chat["id"])
                    delete_chat(chat["id"])
                    remaining = list_chats()
                    if deleting_current:
                        if remaining:
                            st.session_state.current_chat_id = remaining[0]["id"]
                        else:
                            st.session_state.current_chat_id = create_new_chat()
                    st.rerun()

    st.divider()
    st.header("Agent 역할")
    sidebar_chat_data = load_chat(st.session_state.current_chat_id)
    current_agent_role = sidebar_chat_data.get("agent_role", "")
    with st.form(key=f"agent_role_form_{st.session_state.current_chat_id}"):
        edited_agent_role = st.text_input(
            "이 대화의 역할",
            value=current_agent_role,
            placeholder="예: 너는 서비스 기획자야. 기획자 관점으로 답변해줘.",
        )
        col_apply, col_clear = st.columns(2)
        with col_apply:
            apply_clicked = st.form_submit_button("적용하기", use_container_width=True)
        with col_clear:
            clear_clicked = st.form_submit_button("역할 비우기", use_container_width=True)

    if apply_clicked:
        update_chat_agent_role(st.session_state.current_chat_id, edited_agent_role)
        st.success("이 대화의 Agent 역할이 적용되었습니다.")
        st.rerun()

    if clear_clicked:
        update_chat_agent_role(st.session_state.current_chat_id, "")
        st.success("Agent 역할을 비웠습니다.")
        st.rerun()

    st.caption("예시: '너는 PM이야', '너는 마케팅 카피라이터야', '너는 여행 플래너야'")

    st.divider()
    st.header("답변 설정")

    _model_opts = ["gpt-4o-mini", "gpt-4.1-mini", "gpt-4.1", "gpt-5.4"]
    if st.session_state.model_name not in _model_opts:
        st.session_state.model_name = "gpt-4.1-mini"
    st.session_state.model_name = st.selectbox(
        "모델", _model_opts,
        index=_model_opts.index(st.session_state.model_name)
    )

    _len_opts = ["짧게", "보통", "자세히"]
    if st.session_state.answer_length not in _len_opts:
        st.session_state.answer_length = "보통"
    st.session_state.answer_length = st.selectbox(
        "답변 길이", _len_opts,
        index=_len_opts.index(st.session_state.answer_length)
    )

    st.divider()
    st.header("검색 설정")

    st.session_state.use_web_search = st.toggle(
        "검색 사용",
        value=st.session_state.use_web_search
    )

    st.session_state.auto_search_only = st.toggle(
        "검색 필요 질문만 자동검색",
        value=st.session_state.auto_search_only
    )

    st.session_state.show_search_images = st.toggle(
        "네이버 이미지 보기",
        value=st.session_state.show_search_images
    )

    st.session_state.show_web_sources = st.toggle(
        "OpenAI 웹 출처 보기",
        value=st.session_state.show_web_sources
    )

    st.caption("국내 장소/맛집/이미지는 네이버, 최신 뉴스/웹정보는 OpenAI 웹검색을 함께 사용합니다. 질문 성격에 따라 로컬·뉴스·웹문서·이미지를 자동 조합합니다.")

    st.divider()
    st.header("📚 RAG 문서 Q&A")

    st.session_state.rag_enabled = st.toggle(
        "RAG 모드 (문서 기반 답변)",
        value=st.session_state.rag_enabled,
    )
    st.caption("켜면 인덱싱된 문서에서 관련 청크만 찾아 답변에 활용합니다.")

    rag_username = st.session_state.get("username", "guest")
    rag_docs = list_rag_docs(rag_username)

    rag_upload = st.file_uploader(
        "문서 인덱싱 (PDF / DOCX / TXT)",
        type=["pdf", "docx", "txt"],
        key="rag_uploader",
        label_visibility="collapsed",
    )

    if rag_upload and not st.session_state.rag_indexing:
        if st.button("인덱싱 시작", use_container_width=True):
            st.session_state.rag_indexing = True
            with st.spinner(f"'{rag_upload.name}' 임베딩 중..."):
                ext = rag_upload.name.split(".")[-1].lower()
                if ext == "pdf":
                    raw_text = read_pdf(rag_upload)
                elif ext == "docx":
                    raw_text = read_docx(rag_upload)
                else:
                    raw_text = read_txt(rag_upload)

                if raw_text and not raw_text.startswith("["):
                    doc_id = hashlib.md5(f"{rag_username}{rag_upload.name}".encode()).hexdigest()
                    n = index_document(rag_username, doc_id, rag_upload.name, raw_text)
                    st.success(f"완료: {n}개 청크 저장됨")
                else:
                    st.error("문서 읽기 실패")
            st.session_state.rag_indexing = False
            st.rerun()

    if rag_docs:
        st.caption(f"인덱싱된 문서 {len(rag_docs)}건")
        for doc in rag_docs:
            col_d1, col_d2 = st.columns([5, 1])
            with col_d1:
                st.caption(f"📄 {doc['filename']} ({doc['chunks']}청크)")
            with col_d2:
                if st.button("🗑", key=f"rag_del_{doc['doc_id']}"):
                    delete_rag_doc(rag_username, doc["doc_id"])
                    st.rerun()
    else:
        st.caption("인덱싱된 문서 없음")

# ---------------------------------
# 현재 대화 로드
# ---------------------------------
current_data = load_chat(st.session_state.current_chat_id)
messages = current_data["messages"]
current_agent_role = current_data.get("agent_role", "").strip()

if current_agent_role:
    st.info(f"🧠 현재 대화 Agent 역할: {current_agent_role}")

# ---------------------------------
# 파일 첨부 안내 (채팅창 첨부 사용)
# ---------------------------------
st.caption("📎 파일/스크린샷 첨부는 아래 대화 입력창(+)에서 해주세요.")

active_files = st.session_state.uploaded_files_cache

file_context = ""
image_inputs = []

if active_files:
    st.success(f"{len(active_files)}개 파일 업로드됨")

    for f in active_files:
        ext = f.name.split(".")[-1].lower()
        st.write("첨부됨:", f.name)

        try:
            if ext == "pdf":
                text = read_pdf(f)
                file_context += f"\n\n[PDF: {f.name}]\n{text}"

            elif ext in ["xlsx", "xls"]:
                excel_text, previews = read_excel(f)
                file_context += f"\n\n[EXCEL: {f.name}]\n{excel_text}"
                for sheet_name, df in previews:
                    with st.expander(f"미리보기: {f.name} / {sheet_name}", expanded=False):
                        st.dataframe(df, use_container_width=True)

            elif ext == "csv":
                csv_text, preview_df = read_csv(f)
                file_context += f"\n\n[CSV: {f.name}]\n{csv_text}"
                if preview_df is not None:
                    with st.expander(f"미리보기: {f.name}", expanded=False):
                        st.dataframe(preview_df, use_container_width=True)

            elif ext == "pptx":
                text = read_ppt(f)
                file_context += f"\n\n[PPTX: {f.name}]\n{text}"

            elif ext == "docx":
                text = read_docx(f)
                file_context += f"\n\n[DOCX: {f.name}]\n{text}"

            elif ext == "txt":
                text = read_txt(f)
                file_context += f"\n\n[TXT: {f.name}]\n{text}"

            elif ext in ["png", "jpg", "jpeg", "webp"]:
                with st.expander(f"이미지 미리보기: {f.name}", expanded=False):
                    st.image(f, caption=f"{f.name} 원본", use_container_width=True)

                file_context += f"""
[이미지 파일: {f.name}]
이 이미지는 사용자가 첨부한 원본 이미지입니다.
OCR 전처리 텍스트는 제공하지 않으니, 필요한 경우 이미지 자체를 직접 분석하세요.
여권, 비자, 신분증, 문서 이미지, 캡처 화면일 수 있으므로
이름, 여권번호, 국적, 생년월일, 발급일, 만료일, 비자 종류, 체류기간 등의 정보가 보이면 정리하세요.
"""

                image_inputs.append({
                    "type": "input_image",
                    "image_url": f"data:{f.type};base64,{image_to_base64(f)}"
                })

        except Exception as e:
            st.error(f"{f.name} 처리 중 오류: {e}")

    if st.button("첨부 파일 비우기"):
        st.session_state.uploaded_files_cache = []
        st.session_state.last_paste_signature = ""
        st.rerun()
else:
    st.info("업로드된 파일 없음")

with st.expander("첨부 데이터 확인", expanded=False):
    st.write("file_context 길이:", len(file_context))
    st.write("image_inputs 개수:", len(image_inputs))

# ---------------------------------
# 이전 대화 출력
# ---------------------------------
for msg in messages:
    with st.chat_message(msg["role"], avatar=get_chat_avatar(msg["role"])):
        st.write(msg["content"])

# ---------------------------------
# 마지막 HTML 미리보기 재표시
# ---------------------------------
if st.session_state.last_generated_images:
    render_generated_images(st.session_state.last_generated_images)

if st.session_state.last_preview_html:
    render_preview_panel(
        st.session_state.last_preview_html,
        st.session_state.last_preview_blocks,
        key_prefix=f"chat_{st.session_state.current_chat_id}_last"
    )

# ---------------------------------
# 사용자 입력 (채팅창 첨부 지원)
# ---------------------------------
st.caption("💡 채팅창에 포커스를 두고 Ctrl+V로 캡처 이미지를 붙여넣으면 첨부되고, Enter로 바로 전송됩니다.")

def mount_clipboard_image_bridge(chat_key: str):
    st.html(
        f"""
        <script>
        (() => {{
          const ROOT_CLASS = "st-key-{chat_key}";
          if (window.__clipboard_bridge_installed__ === ROOT_CLASS) return;
          window.__clipboard_bridge_installed__ = ROOT_CLASS;

          function findRoot() {{
            return document.querySelector("." + ROOT_CLASS);
          }}

          function findTextarea(root) {{
            if (!root) return null;
            return root.querySelector("textarea");
          }}

          function findFileInput(root) {{
            if (!root) return null;
            return root.querySelector('input[type="file"]');
          }}

          function showBadge(root, text) {{
            if (!root) return;
            let badge = root.querySelector(".clipboard-paste-badge");
            if (!badge) {{
              badge = document.createElement("div");
              badge.className = "clipboard-paste-badge";
              badge.style.cssText = `
                margin-top: 6px;
                font-size: 12px;
                color: #2563eb;
                background: #eff6ff;
                border: 1px solid #bfdbfe;
                border-radius: 999px;
                padding: 4px 10px;
                display: inline-block;
              `;
              root.appendChild(badge);
            }}
            badge.textContent = text;
            clearTimeout(badge._timer);
            badge._timer = setTimeout(() => {{
              if (badge) badge.remove();
            }}, 2500);
          }}

          async function attachClipboardImage(blob) {{
            const root = findRoot();
            const fileInput = findFileInput(root);
            const textarea = findTextarea(root);

            if (!root || !fileInput) return;

            const fileName = `clipboard_${{Date.now()}}.png`;
            const file = new File([blob], fileName, {{
              type: blob.type || "image/png",
              lastModified: Date.now()
            }});

            const dt = new DataTransfer();

            const existingFiles = Array.from(fileInput.files || []);
            for (const f of existingFiles) {{
              dt.items.add(f);
            }}
            dt.items.add(file);

            fileInput.files = dt.files;
            fileInput.dispatchEvent(new Event("change", {{ bubbles: true }}));
            fileInput.dispatchEvent(new Event("input", {{ bubbles: true }}));

            showBadge(root, "📎 클립보드 이미지 첨부됨 · Enter로 전송");
            if (textarea) {{
              textarea.focus();
            }}
          }}

          document.addEventListener("paste", async (event) => {{
            const root = findRoot();
            const textarea = findTextarea(root);
            if (!root || !textarea) return;

            const active = document.activeElement;
            const isChatFocused = active === textarea || root.contains(active);
            if (!isChatFocused) return;

            const items = Array.from(event.clipboardData?.items || []);
            const imageItem = items.find(item => item.type && item.type.startsWith("image/"));
            if (!imageItem) return;

            const blob = imageItem.getAsFile();
            if (!blob) return;

            event.preventDefault();
            await attachClipboardImage(blob);
          }}, true);
        }})();
        </script>
        """,
        unsafe_allow_javascript=True,
    )

def mount_toolbar_toggle():
    st.html(
        """
        <script>
        (() => {
            if (window.__toolbar_toggle_installed__) return;
            window.__toolbar_toggle_installed__ = true;

            const STORAGE_KEY = 'hazel_toolbar_visible';

            function isMobile() {
                return window.innerWidth <= 768;
            }

            function isVisible() {
                return localStorage.getItem(STORAGE_KEY) === 'true';
            }

            function applyState(visible) {
                if (isMobile()) return; // 모바일은 항상 표시
                const toolbar = document.querySelector('[data-testid="stToolbar"]');
                if (toolbar) {
                    toolbar.style.opacity        = visible ? '1' : '0';
                    toolbar.style.pointerEvents  = visible ? 'auto' : 'none';
                    toolbar.style.transition     = 'opacity 0.25s ease';
                }
                // 사이드바 토글은 항상 보이도록 유지
                ['[data-testid="collapsedControl"]','[data-testid="stSidebarCollapsedControl"]']
                    .forEach(sel => {
                        const el = document.querySelector(sel);
                        if (el) { el.style.opacity = '1'; el.style.pointerEvents = 'auto'; }
                    });
                localStorage.setItem(STORAGE_KEY, visible ? 'true' : 'false');
                const btn = document.getElementById('__toolbar_toggle_btn__');
                if (btn) btn.title = visible ? '툴바 숨기기' : '툴바 보이기';
            }

            function createBtn() {
                if (document.getElementById('__toolbar_toggle_btn__')) return;
                const btn = document.createElement('button');
                btn.id = '__toolbar_toggle_btn__';
                btn.innerHTML = '⋯';
                Object.assign(btn.style, {
                    position      : 'fixed',
                    top           : '6px',
                    right         : '6px',
                    zIndex        : '99999',
                    width         : '22px',
                    height        : '22px',
                    background    : 'rgba(245,240,232,0.55)',
                    border        : '1px solid rgba(200,185,170,0.5)',
                    borderRadius  : '6px',
                    cursor        : 'pointer',
                    fontSize      : '14px',
                    lineHeight    : '1',
                    color         : '#8a7560',
                    display       : 'flex',
                    alignItems    : 'center',
                    justifyContent: 'center',
                    padding       : '0',
                    opacity       : '0.35',
                    transition    : 'opacity 0.2s, background 0.2s',
                    backdropFilter: 'blur(4px)',
                });
                btn.onmouseenter = () => { btn.style.opacity = '1'; btn.style.background = 'rgba(245,240,232,0.95)'; };
                btn.onmouseleave = () => { btn.style.opacity = '0.35'; btn.style.background = 'rgba(245,240,232,0.55)'; };
                btn.onclick = () => applyState(!isVisible());
                document.body.appendChild(btn);
            }

            function init() {
                const toolbar = document.querySelector('[data-testid="stToolbar"]');
                if (toolbar) {
                    createBtn();
                    applyState(isVisible());
                } else {
                    setTimeout(init, 150);
                }
            }
            init();

            // Streamlit 리렌더 후에도 상태 재적용
            const observer = new MutationObserver(() => {
                if (!isMobile()) {
                    const toolbar = document.querySelector('[data-testid="stToolbar"]');
                    if (toolbar && (toolbar.style.opacity === '' || toolbar.style.opacity === '1') && !isVisible()) {
                        applyState(false);
                    }
                    createBtn();
                }
            });
            observer.observe(document.body, { childList: true, subtree: true });
        })();
        </script>
        """,
        unsafe_allow_javascript=True,
    )

mount_toolbar_toggle()

if st.session_state.is_generating:
    if st.button("⏹ 응답 멈춤", use_container_width=True):
        st.session_state.stop_generation = True
        st.rerun()

chat_input_file_types = [
    "pdf", "xlsx", "xls", "csv",
    "pptx", "docx", "txt",
    "png", "jpg", "jpeg", "webp"
]

user_input = ""
chat_input_files = []
legacy_chat_uploader_files = []


try:
    chat_payload = st.chat_input(
        "메시지를 입력하세요 (+ 버튼으로 파일 첨부 / Ctrl+V 이미지 붙여넣기)",
        accept_file="multiple",
        file_type=chat_input_file_types,
        key="main_chat_input_with_file",
    )
except Exception:
    chat_payload = st.chat_input("메시지를 입력하세요", key="main_chat_input_fallback")
    legacy_chat_uploader_files = st.file_uploader(
        "파일 첨부(호환 모드)",
        type=chat_input_file_types,
        accept_multiple_files=True,
        key="legacy_chat_uploader",
        label_visibility="collapsed",
    ) or []

mount_clipboard_image_bridge("main_chat_input_with_file")

if isinstance(chat_payload, str):
    user_input = chat_payload.strip()
elif chat_payload is not None:
    payload_text = (getattr(chat_payload, "text", "") or "").strip()
    if payload_text:
        user_input = payload_text
    payload_files = list(getattr(chat_payload, "files", []) or [])
    if payload_files:
        chat_input_files.extend(payload_files)

if legacy_chat_uploader_files:
    chat_input_files.extend(list(legacy_chat_uploader_files))

def _files_signature(files):
    if not files:
        return ""
    parts = []
    for f in files:
        try:
            raw = f.getvalue()
        except Exception:
            f.seek(0)
            raw = f.read()
            f.seek(0)
        digest = hashlib.sha1(raw).hexdigest()
        parts.append(f"{getattr(f, 'name', 'file')}:{digest}")
    return "|".join(parts)

# 중복 제거
unique_files = []
seen_signatures = set()
for f in chat_input_files:
    sig = _files_signature([f])
    if sig and sig not in seen_signatures:
        seen_signatures.add(sig)
        unique_files.append(f)
chat_input_files = unique_files

submitted_text = (user_input or "").strip()
has_chat_submission = bool(submitted_text) or bool(chat_input_files)

if has_chat_submission:
    if chat_input_files:
        st.session_state.uploaded_files_cache = chat_input_files
        if not submitted_text:
            st.info("클립보드/파일 첨부가 감지되었습니다. 이미지 내용을 분석해드릴게요.")

    if not submitted_text:
        submitted_text = "첨부한 파일(이미지/문서)을 분석해줘."

    chat_id = st.session_state.current_chat_id

    messages.append({"role": "user", "content": submitted_text})

    if current_data.get("title") in ["새 대화", "제목 없음"]:
        new_title = make_title_from_messages(messages)
        current_data["title"] = new_title
        update_chat_title(chat_id, new_title)

    append_message(chat_id, "user", submitted_text)

    with st.chat_message("user", avatar=get_chat_avatar("user")):
        st.write(submitted_text)
        if chat_input_files:
            for f in chat_input_files:
                file_name = getattr(f, "name", "첨부 파일")
                ext = file_name.split(".")[-1].lower() if "." in file_name else ""
                if ext in ["png", "jpg", "jpeg", "webp"]:
                    try:
                        f.seek(0)
                    except Exception:
                        pass
                    st.image(f, caption=file_name, width=220)
                else:
                    st.caption(f"첨부 파일: {file_name}")

    with st.chat_message("assistant", avatar=get_chat_avatar("assistant")):
        placeholder = st.empty()
        full_text = ""
        st.session_state.is_generating = True
        st.session_state.stop_generation = False
        runtime_state = get_default_runtime_state()
        naver_results_map = runtime_state["naver_results_map"]
        openai_web_sources = runtime_state["openai_web_sources"]
        generated_images = runtime_state["generated_images"]
        search_plan = runtime_state["search_plan"]
        do_search = runtime_state["do_search"]
        usage_input_tokens = 0
        usage_output_tokens = 0
        rag_chunks_used = []

        try:
            history_for_model = []
            for msg in messages[:-1]:
                history_for_model.append({
                    "role": msg["role"],
                    "content": msg["content"]
                })

            if is_image_generation_request(submitted_text) and not should_prioritize_code_preview(submitted_text, effective_files):
                image_prompt = get_image_generation_prompt(submitted_text)
                image_spec = build_image_generation_spec(submitted_text, image_prompt)
                placeholder.info("🎨 이미지 생성 중입니다... 잠시만 기다려주세요.")
                generated_images = generate_openai_image(
                    image_spec["prompt"],
                    size=image_spec["size"],
                    quality=image_spec["quality"]
                )
                st.session_state.last_generated_images = generated_images
                st.session_state.last_generated_prompt = image_spec["prompt"]

                if generated_images:
                    full_text = f"""요청한 이미지 생성이 완료되었습니다. 아래에서 결과를 확인하고 다운로드할 수 있어요.

적용 모드: {image_spec["mode"]}
프롬프트: {image_prompt}"""
                    placeholder.markdown(full_text)
                else:
                    full_text = "이미지 생성 결과를 받지 못했습니다. 프롬프트를 조금 더 구체적으로 적어 주세요."
                    placeholder.warning(full_text)
            else:
                st.session_state.last_generated_images = []
                st.session_state.last_generated_prompt = ""

                if st.session_state.use_web_search:
                    if st.session_state.auto_search_only:
                        do_search = should_search_web(submitted_text)
                    else:
                        do_search = True

                search_plan = build_search_plan(submitted_text) if do_search else {"mode_labels": ["일반 응답"]}

                # RAG 검색
                rag_context = ""
                rag_chunks_used = []
                if st.session_state.rag_enabled:
                    rag_chunks_used = search_rag_chunks(
                        st.session_state.get("username", "guest"),
                        submitted_text,
                        top_k=RAG_TOP_K,
                    )
                    if rag_chunks_used:
                        rag_context = "\n\n".join(
                            f"[{c['filename']} / 청크{c['chunk_index']}]\n{c['text']}"
                            for c in rag_chunks_used
                        )

                user_text = f"""사용자 질문:
{submitted_text}

첨부 파일 내용:
{file_context if file_context else "첨부된 파일 없음"}
"""
                if rag_context:
                    user_text += f"""
RAG 문서 관련 내용 (가장 유사한 청크):
{rag_context}
"""

                if do_search and search_plan.get("use_local"):
                    naver_results_map["local"] = naver_search(submitted_text, search_type="local", display=5)
                    user_text += f"""

네이버 장소 검색 결과:
{format_naver_search_results(naver_results_map['local'], search_type='local')}
"""

                if do_search and search_plan.get("use_naver_news"):
                    naver_results_map["news"] = naver_search(submitted_text, search_type="news", display=4)
                    user_text += f"""

네이버 뉴스 검색 결과:
{format_naver_search_results(naver_results_map['news'], search_type='news')}
"""

                if do_search and search_plan.get("use_naver_web"):
                    naver_results_map["web"] = naver_search(submitted_text, search_type="webkr", display=4)
                    user_text += f"""

네이버 웹문서 검색 결과:
{format_naver_search_results(naver_results_map['web'], search_type='webkr')}
"""

                if do_search and st.session_state.show_search_images and search_plan.get("wants_images"):
                    naver_results_map["image"] = naver_image_search(make_image_search_query(submitted_text), display=6)
                    user_text += f"""

네이버 이미지 검색 결과:
{format_image_search_results(naver_results_map['image'])}
"""

                user_content = [
                    {
                        "type": "input_text",
                        "text": user_text
                    }
                ]

                if image_inputs:
                    user_content.extend(image_inputs)

                usage_input_tokens = 0
                usage_output_tokens = 0

                if do_search and search_plan.get("use_openai_web"):
                    full_text, openai_web_sources, (usage_input_tokens, usage_output_tokens) = run_openai_web_search(
                        model_name=st.session_state.model_name,
                        instructions=build_system_prompt(
                            st.session_state.answer_length,
                            current_agent_role
                        ),
                        history_for_model=history_for_model,
                        user_content=user_content
                    )
                    placeholder.markdown(full_text)
                else:
                    stream = client.responses.create(
                        model=st.session_state.model_name,
                        input=[
                            {
                                "role": "system",
                                "content": build_system_prompt(
                                    st.session_state.answer_length,
                                    current_agent_role
                                )
                            },
                            *history_for_model,
                            {"role": "user", "content": user_content}
                        ],
                        stream=True
                    )

                    for event in stream:
                        if st.session_state.stop_generation:
                            full_text += "\n\n(사용자 요청으로 응답 생성을 중단했습니다.)"
                            break
                        if event.type == "response.output_text.delta":
                            full_text += event.delta
                            placeholder.markdown(full_text + "▌")
                        elif event.type == "response.completed":
                            try:
                                usage_input_tokens = event.response.usage.input_tokens or 0
                                usage_output_tokens = event.response.usage.output_tokens or 0
                            except Exception:
                                pass
                            break

                    placeholder.markdown(full_text)

        except Exception as e:
            full_text = f"오류가 발생했습니다: {e}"
            placeholder.error(full_text)
        finally:
            st.session_state.is_generating = False
            st.session_state.stop_generation = False

        messages.append({"role": "assistant", "content": full_text})
        append_message(chat_id, "assistant", full_text)

        if usage_input_tokens or usage_output_tokens:
            st.markdown(
                f"<div style='font-size:0.8rem;color:#8a7560;margin-top:4px'>"
                f"{calc_usage_display(st.session_state.model_name, usage_input_tokens, usage_output_tokens)}"
                f"</div>",
                unsafe_allow_html=True
            )

        if rag_chunks_used:
            with st.expander(f"📚 RAG 참조 청크 {len(rag_chunks_used)}개", expanded=False):
                for c in rag_chunks_used:
                    st.markdown(f"**{c['filename']} / 청크 {c['chunk_index']}** (유사도: {c.get('score', 0):.3f})")
                    st.caption(c["text"][:300] + ("…" if len(c["text"]) > 300 else ""))

        if generated_images:
            render_generated_images(generated_images)

        if do_search:
            st.markdown(
                format_search_summary(search_plan, naver_results_map, openai_web_sources),
                unsafe_allow_html=True
            )

        if naver_results_map["local"]:
            with st.expander("네이버 장소 검색 결과 보기", expanded=False):
                render_naver_search_results(naver_results_map["local"], search_type="local")

        if naver_results_map["news"]:
            with st.expander("네이버 뉴스 검색 결과 보기", expanded=False):
                render_naver_search_results(naver_results_map["news"], search_type="news")

        if naver_results_map["web"]:
            with st.expander("네이버 웹문서 검색 결과 보기", expanded=False):
                render_naver_search_results(naver_results_map["web"], search_type="webkr")

        valid_image_results = get_valid_image_results(naver_results_map["image"])
        if valid_image_results:
            with st.expander("관련 이미지 보기", expanded=False):
                render_image_results(valid_image_results)

        if openai_web_sources and st.session_state.show_web_sources:
            with st.expander("OpenAI 웹검색 출처 보기", expanded=False):
                render_openai_web_sources(openai_web_sources)

        result_df = try_build_result_dataframe(full_text)
        st.session_state.last_result_df = result_df

        if result_df is not None and not result_df.empty:
            st.subheader("📊 AI 결과 표")
            st.dataframe(result_df, use_container_width=True)

            result_excel = dataframe_to_excel_bytes(result_df, sheet_name="ai_result")
            st.download_button(
                label="📥 AI 결과 Excel 다운로드",
                data=result_excel,
                file_name="ai_result.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key=f"download_ai_excel_{st.session_state.current_chat_id}"
            )

        st.session_state.last_preview_html = None
        st.session_state.last_preview_blocks = {"html": "", "css": "", "js": ""}

        if should_show_preview(submitted_text, full_text):
            preview_html, preview_blocks = build_preview_html_from_response(full_text)

            if preview_html:
                st.session_state.last_preview_html = preview_html
                st.session_state.last_preview_blocks = preview_blocks
                render_preview_panel(
                    preview_html,
                    preview_blocks,
                    key_prefix=f"chat_{chat_id}_live"
                )
            else:
                if "```css" in full_text.lower() and "```html" not in full_text.lower():
                    st.info("CSS 코드만 있어서 미리보기는 생략했습니다. HTML 코드까지 같이 있으면 바로 렌더됩니다.")
