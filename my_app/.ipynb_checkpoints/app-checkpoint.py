import os
import json
import re
import base64
from datetime import datetime

import streamlit as st
import pandas as pd

from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from openai import OpenAI

# ----------------------------
# 기본 설정
# ----------------------------
st.set_page_config(
    page_title="내 AI 챗봇",
    page_icon="🤖",
    layout="wide"
)

st.title("🤖 AI 챗봇")

CHAT_DIR = "chats"
os.makedirs(CHAT_DIR, exist_ok=True)

# ----------------------------
# OpenAI
# ----------------------------
api_key = os.getenv("OPENAI_API_KEY")

if not api_key:
    st.error("OPENAI_API_KEY 환경변수가 없습니다.")
    st.stop()

client = OpenAI(api_key=api_key)

# ----------------------------
# 파일 읽기 함수
# ----------------------------
def read_pdf(file):
    try:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        return f"[PDF 읽기 실패: {e}]"


def read_excel(file):
    try:
        excel = pd.ExcelFile(file)
        text_parts = []
        preview_dfs = []

        for sheet_name in excel.sheet_names:
            df = pd.read_excel(excel, sheet_name=sheet_name)
            preview_dfs.append((sheet_name, df.head(20)))
            text_parts.append(f"[시트: {sheet_name}]")
            text_parts.append(df.head(50).to_string(index=False))

        return "\n\n".join(text_parts), preview_dfs
    except Exception as e:
        return f"[Excel 읽기 실패: {e}]", []


def read_csv(file):
    try:
        df = pd.read_csv(file)
    except Exception:
        try:
            file.seek(0)
            df = pd.read_csv(file, encoding="cp949")
        except Exception as e:
            return f"[CSV 읽기 실패: {e}]", None

    return df.head(50).to_string(index=False), df.head(20)


def read_ppt(file):
    try:
        prs = Presentation(file)
        text = ""
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    value = shape.text.strip()
                    if value:
                        slide_texts.append(value)
            text += f"\n[슬라이드 {i}]\n" + "\n".join(slide_texts) + "\n"
        return text.strip()
    except Exception as e:
        return f"[PPT 읽기 실패: {e}]"


def read_docx(file):
    try:
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        return f"[DOCX 읽기 실패: {e}]"


def read_txt(file):
    raw = file.getvalue()
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(enc)
        except Exception:
            pass
    return raw.decode("utf-8", errors="ignore")


def image_to_base64(file):
    return base64.b64encode(file.getvalue()).decode()


# ----------------------------
# 대화 저장 함수
# ----------------------------
def chat_path(chat_id):
    return os.path.join(CHAT_DIR, f"{chat_id}.json")


def create_new_chat():
    chat_id = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    data = {
        "title": "새 대화",
        "messages": [
            {"role": "assistant", "content": "안녕하세요! 무엇을 도와드릴까요?"}
        ]
    }
    with open(chat_path(chat_id), "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return chat_id


def load_chat(chat_id):
    path = chat_path(chat_id)
    if not os.path.exists(path):
        return {
            "title": "새 대화",
            "messages": [
                {"role": "assistant", "content": "안녕하세요! 무엇을 도와드릴까요?"}
            ]
        }

    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def save_chat(chat_id, data):
    with open(chat_path(chat_id), "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def list_chats():
    chats = []
    for file in os.listdir(CHAT_DIR):
        if file.endswith(".json"):
            chat_id = file.replace(".json", "")
            try:
                data = load_chat(chat_id)
                chats.append({"id": chat_id, "title": data.get("title", "제목 없음")})
            except Exception:
                pass

    chats.sort(reverse=True, key=lambda x: x["id"])
    return chats


# ----------------------------
# 세션 초기화
# ----------------------------
if "chat_id" not in st.session_state:
    chats = list_chats()
    if chats:
        st.session_state.chat_id = chats[0]["id"]
    else:
        st.session_state.chat_id = create_new_chat()

if "uploaded_files_cache" not in st.session_state:
    st.session_state.uploaded_files_cache = []

# ----------------------------
# 사이드바
# ----------------------------
with st.sidebar:
    if st.button("➕ 새 대화", use_container_width=True):
        st.session_state.chat_id = create_new_chat()
        st.session_state.uploaded_files_cache = []
        st.rerun()

    st.divider()
    st.subheader("대화 목록")

    for chat in list_chats():
        col1, col2 = st.columns([4, 1])

        with col1:
            if st.button(chat["title"], key=f"open_{chat['id']}", use_container_width=True):
                st.session_state.chat_id = chat["id"]
                st.session_state.uploaded_files_cache = []
                st.rerun()

        with col2:
            if st.button("🗑", key=f"del_{chat['id']}", use_container_width=True):
                path = chat_path(chat["id"])
                if os.path.exists(path):
                    os.remove(path)

                remain = list_chats()
                if remain:
                    st.session_state.chat_id = remain[0]["id"]
                else:
                    st.session_state.chat_id = create_new_chat()

                st.session_state.uploaded_files_cache = []
                st.rerun()

    st.divider()
    st.subheader("답변 설정")

    model_name = st.selectbox(
        "모델",
        ["gpt-4o-mini", "gpt-4.1-mini", "gpt-4.1"],
        index=1
    )

    answer_length = st.selectbox(
        "답변 길이",
        ["짧게", "보통", "자세히"],
        index=1
    )

# ----------------------------
# 프롬프트 생성
# ----------------------------
def build_prompt():
    if answer_length == "짧게":
        length = "답변은 2~3문장으로 간단히 설명한다."
    elif answer_length == "보통":
        length = "답변은 3~6문장 정도로 설명한다."
    else:
        length = "답변은 충분히 자세히 설명하고 예시도 포함한다."

    return f"""
너는 친절한 한국어 AI 챗봇이다.
항상 한국어로 답변한다.
모르는 내용은 추측하지 말고 모른다고 말한다.
사용자가 파일을 첨부한 경우 첨부 내용도 함께 참고한다.
{length}
사용자가 이해할 수 있도록 설명한다.
"""

system_prompt = build_prompt()

# ----------------------------
# 파일 업로드
# ----------------------------
st.subheader("📎 파일 첨부")

uploaded_files = st.file_uploader(
    "여기에 파일을 드래그하거나 클릭해서 선택하세요",
    type=[
        "pdf", "xlsx", "xls", "csv",
        "pptx", "docx", "txt",
        "png", "jpg", "jpeg", "webp"
    ],
    accept_multiple_files=True,
    key="main_file_uploader"
)

st.write("현재 업로드 객체:", uploaded_files)

if uploaded_files is not None and len(uploaded_files) > 0:
    st.session_state.uploaded_files_cache = uploaded_files
    st.success(f"{len(uploaded_files)}개 파일 업로드됨")
else:
    if "uploaded_files_cache" not in st.session_state:
        st.session_state.uploaded_files_cache = []

active_files = st.session_state.uploaded_files_cache

if active_files:
    for f in active_files:
        st.write("첨부됨:", f.name)
else:
    st.info("업로드된 파일 없음")

# ----------------------------
# 대화 불러오기
# ----------------------------
chat_data = load_chat(st.session_state.chat_id)
messages = chat_data["messages"]

# ----------------------------
# 채팅 출력
# ----------------------------
for msg in messages:
    with st.chat_message(msg["role"]):
        st.write(msg["content"])

# ----------------------------
# 사용자 입력
# ----------------------------
user_input = st.chat_input("메시지를 입력하세요")

if user_input:
    messages.append({"role": "user", "content": user_input})

    if chat_data["title"] == "새 대화":
        chat_data["title"] = user_input[:25]

    chat_data["messages"] = messages
    save_chat(st.session_state.chat_id, chat_data)

    with st.chat_message("user"):
        st.write(user_input)

    with st.chat_message("assistant"):
        placeholder = st.empty()
        full_text = ""

        try:
            history_for_model = []
            for msg in messages[:-1]:
                history_for_model.append({
                    "role": msg["role"],
                    "content": msg["content"]
                })

            user_content = [
                {
                    "type": "input_text",
                    "text": f"""사용자 질문:
{user_input}

첨부 파일 내용:
{file_context if file_context else "첨부된 파일 없음"}
"""
                }
            ]

            if image_inputs:
                user_content.extend(image_inputs)

            response = client.responses.create(
                model=model_name,
                input=[
                    {"role": "system", "content": system_prompt},
                    *history_for_model,
                    {"role": "user", "content": user_content}
                ],
                stream=True
            )

            for event in response:
                if event.type == "response.output_text.delta":
                    full_text += event.delta
                    placeholder.markdown(full_text + "▌")
                elif event.type == "response.completed":
                    break

            placeholder.markdown(full_text)

        except Exception as e:
            full_text = f"오류가 발생했습니다: {e}"
            placeholder.error(full_text)

        messages.append({"role": "assistant", "content": full_text})
        chat_data["messages"] = messages
        save_chat(st.session_state.chat_id, chat_data)

        code_blocks = re.findall(r"```(?:\w+)?\n(.*?)```", full_text, re.DOTALL)

        for i, code in enumerate(code_blocks):
            st.code(code.strip())
            st.download_button(
                label=f"📋 코드 저장 {i+1}",
                data=code.strip(),
                file_name=f"code_block_{i+1}.txt",
                mime="text/plain",
                key=f"download_code_{st.session_state.chat_id}_{i}"
            )